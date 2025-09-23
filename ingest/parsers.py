from __future__ import annotations
from pathlib import Path
from typing import Tuple
from pypdf import PdfReader
import os

# Supported image extensions for OCR
IMAGE_EXTS = {".png", ".jpg", ".jpeg"}

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    import pptx  # python-pptx
except Exception:
    pptx = None


def read_text(path: Path) -> Tuple[str, dict]:
    meta: dict = {"source": str(path), "name": path.name}
    if path.suffix.lower() == ".pdf":
        return _parse_pdf_with_auto_upgrade(path, meta)
    if path.suffix.lower() in IMAGE_EXTS:
        return _parse_image_with_docai(path, meta)
        
def _parse_pdf_with_auto_upgrade(path: Path, base_meta: dict) -> Tuple[str, dict]:
    """
    Parse PDF with auto-upgrade to DocAI when local extraction yields poor results
    """
    from core.config import is_docai_enabled
    from core.logging import logger
    
    # Thresholds for auto-upgrade decision
    LOW_TEXT_THRESHOLD = 500  # characters
    EMPTY_PAGES_THRESHOLD = 0.6  # 60% empty pages
    
    try:
        # Step 1: Try local PDF extraction first
        reader = PdfReader(str(path))
        total_pages = len(reader.pages)
        
        page_texts = []
        empty_pages = 0
        
        for page in reader.pages:
            page_text = page.extract_text() or ""
            page_texts.append(page_text)
            
            if not page_text.strip():
                empty_pages += 1
        
        # Join all pages
        local_text = "\n".join(page_texts)
        text_length = len(local_text.strip())
        empty_pages_ratio = empty_pages / total_pages if total_pages > 0 else 0
        
        # Log local extraction results
        logger.debug(f"Local PDF extraction: {path.name} -> {text_length} chars, "
                    f"{empty_pages}/{total_pages} empty pages ({empty_pages_ratio:.1%})")
        
        # Step 2: Decide whether to upgrade to DocAI
        should_upgrade_to_docai = (
            text_length < LOW_TEXT_THRESHOLD or 
            empty_pages_ratio >= EMPTY_PAGES_THRESHOLD
        )
        
        # Manual override check
        force_docai = is_docai_enabled() and settings.use_google_doc_ai
        
        if force_docai or (should_upgrade_to_docai and is_docai_enabled()):
            logger.info(f"PDF parser: DocAI ({'manual override' if force_docai else 'auto-upgrade'})")
            
            try:
                from gcp.docai import parse_pdf_with_docai
                docai_text, docai_meta = parse_pdf_with_docai(path)
                
                # If DocAI extraction successful and has more content, use it
                if docai_text and len(docai_text.strip()) > text_length:
                    logger.info(f"DocAI extraction successful: {len(docai_text)} chars vs {text_length} local chars")
                    return docai_text, {**base_meta, **docai_meta}
                else:
                    logger.warning("DocAI extraction didn't improve results, using local extraction")
                    
            except Exception as e:
                logger.warning(f"DocAI parsing failed, falling back to local: {e}")
        
        # Step 3: Use local extraction (default or fallback)
        if not force_docai:
            logger.info("PDF parser: PdfReader (local)")
        else:
            logger.info("PDF parser: DocAI failed, falling back to local")
        
        return local_text, {
            **base_meta, 
            "parser": "local", 
            "pages": total_pages,
            "empty_pages": empty_pages,
            "text_length": text_length
        }
        
    except Exception as e:
        logger.error(f"PDF parsing failed completely: {e}")
        return "", {**base_meta, "parser": "error", "error": str(e)}

def _parse_image_with_docai(path: Path, base_meta: dict) -> Tuple[str, dict]:
    """
    Parse image using DocAI OCR or return empty with suggestion
    """
    from core.config import settings
    from core.logging import logger
    
    # Check if DocAI is enabled
    use_docai = (
        getattr(settings, "use_google_doc_ai", False) or 
        os.getenv("USE_GOOGLE_DOC_AI", "false").lower() == "true"
    )
    
    if use_docai:
        try:
            from gcp.docai import parse_pdf_with_docai
            # DocAI can handle images too, just pass the image file
            text, docai_meta = parse_pdf_with_docai(path)
            logger.info(f"Image OCR via DocAI: {path.name} -> {len(text)} characters")
            return text, {**base_meta, **docai_meta}
        except Exception as e:
            logger.warning(f"DocAI image OCR failed for {path.name}: {e}")
            return "", {**base_meta, "parser": "docai_error", "error": str(e)}
    else:
        logger.info(f"Image file detected: {path.name} - enable USE_GOOGLE_DOC_AI for OCR")
        return "", {**base_meta, "parser": "no_ocr", "note": "Enable USE_GOOGLE_DOC_AI for OCR on images"}

    if path.suffix.lower() in {".txt", ".md"}:
        return path.read_text(errors="ignore"), meta
    if path.suffix.lower() == ".docx" and docx:
        d = docx.Document(str(path))
        return "\n".join(p.text for p in d.paragraphs), meta
    if path.suffix.lower() == ".pptx" and pptx:
        pres = pptx.Presentation(str(path))
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts), meta
    if path.suffix.lower() in {".csv", ".xlsx"}:
        try:
            import pandas as pd
            df = pd.read_csv(path) if path.suffix.lower() == ".csv" else pd.read_excel(path)
            return df.to_csv(index=False), meta
        except Exception:
            return "", meta
    return "", meta
